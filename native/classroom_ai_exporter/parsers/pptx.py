from __future__ import annotations

from pathlib import Path


def extract_pptx_text(path: str | Path) -> str:
    from pptx import Presentation

    presentation = Presentation(str(path))
    slides: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        slides.append(f"# Slide {index}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slides.append(shape.text.strip())
    return "\n".join(slides)
